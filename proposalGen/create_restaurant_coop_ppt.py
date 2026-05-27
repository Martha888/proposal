#!/usr/bin/env python3
"""
ProposalGen - Create Restaurant Cooperation Dish Proposal PPT

Start from examples/sample.pptx and keep the source deck as much as possible:
- Replace the cover text/image.
- Keep supplier introduction pages mostly unchanged.
- Replace slides 6+ with Hunan-style cooperation dish recommendation pages.
- Use generated lotus-root opportunity dish images and write taste/method notes.

Usage:
    python -B proposalGen/create_restaurant_coop_ppt.py
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import urllib.request
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import quote, urlsplit, urlunsplit

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_PROJECT = REPO_ROOT / "projects" / "menu_research" / "费大厨_20260527_095043"
DEFAULT_TEMPLATE = REPO_ROOT / "examples" / "sample.pptx"

RED = RGBColor(0, 167, 60)
DARK = RGBColor(24, 72, 47)
GREEN = RGBColor(0, 167, 60)
MUTED = RGBColor(74, 116, 85)
CREAM = RGBColor(245, 255, 239)
LIGHT = RGBColor(250, 255, 246)
GOLD = RGBColor(238, 227, 136)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(0, 96, 52)

# Shared title safe area. Generated page content must start at or below
# HEADER_CONTENT_TOP so variable titles/subtitles never collide with body copy.
HEADER_LEFT = 560_000
HEADER_TOP = 220_000
HEADER_TITLE_WIDTH = 7_900_000
HEADER_TITLE_HEIGHT = 560_000
HEADER_SECTION_TOP = 860_000
HEADER_SECTION_WIDTH = 5_600_000
HEADER_SECTION_HEIGHT = 310_000
HEADER_RULE_TOP = 1_250_000
HEADER_RULE_HEIGHT = 24_000
HEADER_CONTENT_TOP = 1_540_000


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Create restaurant cooperation proposal PPT.")
    parser.add_argument("--template", default=str(DEFAULT_TEMPLATE), help="Source PPTX template.")
    parser.add_argument("--project-dir", default=str(DEFAULT_PROJECT), help="Project directory.")
    parser.add_argument(
        "--output",
        default="",
        help="Output PPTX filename inside project-dir. Defaults to <restaurant>_lotus_coop_dish_proposal.pptx.",
    )
    parser.add_argument("--restaurant-name", default="", help="Restaurant name. Defaults to metadata in lotus_root_opportunities.json.")
    parser.add_argument(
        "--cover-dish-name",
        default="",
        help="Optional dish name to use as the cover hero image. Defaults to a Hunan-style hot dish.",
    )
    return parser


def _remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _delete_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def _clear_slide_for_rebuild(slide, *, keep_header: bool = True) -> None:
    for shape in list(slide.shapes):
        keep = False
        if keep_header and shape.shape_type == 13:
            # The source template repeats small top-corner brand/decor images.
            keep = shape.top < 500_000 and shape.height < 1_250_000
        if not keep:
            _remove_shape(shape)


def _add_box(slide, left, top, width, height, fill=LIGHT, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line or RGBColor(238, 224, 202)
    box.line.width = Pt(1)
    return box


def _add_text(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    size: int = 18,
    color=DARK,
    bold: bool = False,
    fill=None,
    line=None,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def _add_multiline(
    slide,
    lines: list[tuple[str, int, RGBColor, bool]],
    left: int,
    top: int,
    width: int,
    height: int,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def _add_picture_fit(slide, image_path: Path, left: int, top: int, width: int, height: int):
    if not image_path.exists():
        _add_box(slide, left, top, width, height, fill=RGBColor(245, 230, 220))
        return _add_text(slide, f"图片缺失\n{image_path.name}", left, top + height // 3, width, height // 3, size=14, color=RED, bold=True)
    # Generated images are square; use square frames when possible to avoid distortion.
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return pic


def _crop_picture_to_fill(pic, image_path: Path) -> None:
    try:
        with Image.open(image_path) as img:
            image_ratio = img.width / img.height
    except Exception:
        return

    frame_ratio = pic.width / pic.height
    pic.crop_left = 0
    pic.crop_right = 0
    pic.crop_top = 0
    pic.crop_bottom = 0
    if image_ratio > frame_ratio:
        crop = (1 - frame_ratio / image_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif image_ratio < frame_ratio:
        crop = (1 - image_ratio / frame_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop


def _add_picture_fill(slide, image_path: Path | None, left: int, top: int, width: int, height: int, *, label: str = ""):
    if not image_path or not image_path.exists():
        _add_box(slide, left, top, width, height, fill=RGBColor(232, 214, 196), radius=False)
        placeholder = label or "图片待补充"
        return _add_text(slide, placeholder, left + 180_000, top + height // 2 - 170_000, width - 360_000, 340_000, size=16, color=RED, bold=True)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    _crop_picture_to_fill(pic, image_path)
    return pic


def _replace_picture_content(slide, pic, image_path: Path, *, crop_to_fill: bool = True) -> bool:
    if not image_path.exists():
        return False
    _image_part, r_id = slide.part.get_or_add_image_part(str(image_path))
    pic._element.blipFill.blip.rEmbed = r_id
    if crop_to_fill:
        _crop_picture_to_fill(pic, image_path)
    return True


def _find_cover_hero_picture(slide, prs: Presentation):
    """Find the existing right-side cover image shape without touching background/logo."""
    slide_area = prs.slide_width * prs.slide_height
    candidates = []
    for shape in slide.shapes:
        if shape.shape_type != 13:
            continue

        area = shape.width * shape.height
        center_x = shape.left + shape.width / 2
        is_full_bleed_background = shape.width > prs.slide_width * 0.85 and shape.height > prs.slide_height * 0.85
        is_small_logo_or_decor = area < slide_area * 0.04
        is_right_side_visual = center_x > prs.slide_width * 0.48 and shape.width > prs.slide_width * 0.22

        if is_full_bleed_background or is_small_logo_or_decor or not is_right_side_visual:
            continue
        candidates.append(shape)

    if not candidates:
        return None
    return max(candidates, key=lambda shape: shape.width * shape.height)


def _select_cover_item(items: list[dict[str, Any]], project_dir: Path, preferred_name: str = "") -> dict[str, Any]:
    if preferred_name:
        preferred = next((item for item in items if item.get("name") == preferred_name), None)
        if preferred and _image_path(project_dir, preferred).exists():
            return preferred

    # The cover should feel like a Hunan stir-fry restaurant: hot, spicy, wok-fired, and table-ready.
    priority_keywords = ("干锅", "小炒", "口味", "剁椒", "擂椒", "藕带", "砂锅")
    scored = []
    for idx, item in enumerate(items):
        if not _image_path(project_dir, item).exists():
            continue
        name = str(item.get("name", ""))
        method = str(item.get("method_brief", ""))
        text = f"{name} {method}"
        score = 0
        for weight, keyword in enumerate(priority_keywords[::-1], start=1):
            if keyword in text:
                score += weight
        if _category_for(item) == "湘味热炒":
            score += 5
        scored.append((score, -idx, item))

    if not scored:
        return items[0]
    return max(scored, key=lambda row: (row[0], row[1]))[2]


def _set_shape_text(shape, text: str, *, size: int | None = None, color: RGBColor | None = None, bold: bool | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        if size:
            run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
        if bold is not None:
            run.font.bold = bold


def _replace_shape_text_preserve_format(shape, replacements: dict[str, str]) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False

    changed = False
    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = paragraph.text
        new_text = paragraph_text
        for old, new in replacements.items():
            new_text = new_text.replace(old, new)
        if new_text == paragraph_text:
            continue

        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = new_text
        changed = True
    return changed


def _replace_text_contains(prs: Presentation, slide_no: int, replacements: dict[str, str]) -> None:
    slide = prs.slides[slide_no - 1]
    for shape in slide.shapes:
        _replace_shape_text_preserve_format(shape, replacements)


def _flavor_for(name: str, method: str) -> str:
    text = f"{name} {method}"
    if any(token in text for token in ("藕泥蒸排骨", "剁椒蒸藕片")):
        return "粉糯鲜香、蒸菜温润、湘味足"
    if "口味凤爪" in text:
        return "香辣软糯、藕块吸汁、下饭解腻"
    if any(token in text for token in ("汤", "筒骨", "排骨汤")):
        return "清甜温润、骨香浓、解辣平衡"
    if any(token in text for token in ("凉拌", "酸辣", "藕尖")):
        return "酸辣脆爽、开胃解腻"
    if any(token in text for token in ("桂花", "糯米")):
        return "桂花甜润、软糯收尾"
    if any(token in text for token in ("藕合", "香酥", "炸")):
        return "外酥里嫩、咸香适口"
    if any(token in text for token in ("干锅", "砂锅", "剁椒", "口味", "小炒")):
        return "香辣下饭、锅气足、湘味浓"
    return "莲藕清甜、口感脆爽、适合湘菜桌搭配"


def _category_for(item: dict[str, Any]) -> str:
    name = str(item.get("name", ""))
    method = str(item.get("method_brief", ""))
    if any(token in name for token in ("汤", "筒骨", "排骨汤", "蒸", "砂锅", "焖")):
        return "汤品蒸菜"
    if any(token in name for token in ("凉拌", "酸辣藕尖")):
        return "开胃凉菜"
    if any(token in name for token in ("桂花", "糯米")):
        return "甜品收尾"
    if any(token in name for token in ("藕合", "香酥", "藕夹")):
        return "香酥小吃"
    if "汤" in method and not any(token in name for token in ("口味凤爪", "干锅", "小炒", "擂椒", "藕带")):
        return "汤品蒸菜"
    return "湘味热炒"


def _short(text: str, limit: int = 96) -> str:
    text = re.sub(r"\s+", "", text or "")
    return text if len(text) <= limit else text[:limit] + "..."


def _method_steps(method: str) -> list[str]:
    text = re.sub(r"\s+", "", method or "")
    parts = [part for part in re.split(r"[。；;]", text) if part]
    if not parts:
        return ["方法待细化", "结合门店 SOP 打样", "出餐前复核口味"]

    if len(parts) >= 3:
        steps = [parts[0], parts[1], "".join(parts[2:])]
    elif len(parts) == 2:
        first_split = re.split(r"[，,]", parts[0], maxsplit=1)
        if len(first_split) == 2 and all(first_split):
            steps = [first_split[0], first_split[1], parts[1]]
        else:
            steps = [parts[0], parts[1], "出餐前复核口味与装盘"]
    else:
        chunks = re.split(r"[，,]", parts[0])
        chunks = [chunk for chunk in chunks if chunk]
        if len(chunks) >= 3:
            steps = [chunks[0], chunks[1], "".join(chunks[2:])]
        elif len(chunks) == 2:
            steps = [chunks[0], chunks[1], "出餐前复核口味与装盘"]
        else:
            steps = [parts[0], "结合门店 SOP 打样", "出餐前复核口味与装盘"]

    labels = ["第一步", "第二步", "第三步"]
    return [f"{label}：{_short(step, 34)}" for label, step in zip(labels, steps)]


def _add_method_steps(slide, method: str, left: int, top: int, width: int, height: int, *, color, size: int = 8):
    lines = [(step, size, color, False) for step in _method_steps(method)]
    return _add_multiline(slide, lines, left, top, width, height)


def _image_path(project_dir: Path, item: dict[str, Any]) -> Path:
    generated = item.get("generated_image") or {}
    rel = generated.get("path") or generated.get("url") or ""
    return project_dir / rel


def _support_image_candidates(project_dir: Path) -> list[Path]:
    support_dir = project_dir / "support_images"
    return sorted(support_dir.glob("support_*.png"))


def _support_image_path(project_dir: Path, support_index: int) -> Path | None:
    candidates = _support_image_candidates(project_dir)
    if not candidates:
        return None
    return candidates[(support_index - 1) % len(candidates)]


def _load_opportunities(project_dir: Path) -> list[dict[str, Any]]:
    return _load_opportunity_data(project_dir).get("opportunities", [])


def _load_opportunity_data(project_dir: Path) -> dict[str, Any]:
    return json.loads((project_dir / "lotus_root_opportunities.json").read_text(encoding="utf-8"))


def _restaurant_name_from_project(project_dir: Path, data: dict[str, Any], override: str = "") -> str:
    if override.strip():
        return override.strip()
    name = str(data.get("restaurant") or "").strip()
    if name:
        return name
    structured_menu = project_dir / "structured_menu.json"
    if structured_menu.exists():
        menu_data = json.loads(structured_menu.read_text(encoding="utf-8"))
        name = str(menu_data.get("restaurant") or "").strip()
        if name:
            return name
    return project_dir.name.split("_")[0] or "餐厅"


def _safe_filename(value: str, limit: int = 42) -> str:
    cleaned = re.sub(r"[\\/:*?\"<>|]+", "_", value.strip())
    cleaned = re.sub(r"\s+", "_", cleaned)
    return (cleaned or "dish")[:limit]


def _url_to_request_url(url: str) -> str:
    parts = urlsplit(url)
    path = quote(parts.path, safe="/%")
    query = quote(parts.query, safe="=&?/%")
    return urlunsplit((parts.scheme, parts.netloc, path, query, parts.fragment))


def _load_menu_photo_index(project_dir: Path) -> dict[str, dict[str, Any]]:
    menu_path = project_dir / "structured_menu.json"
    if not menu_path.exists():
        return {}
    data = json.loads(menu_path.read_text(encoding="utf-8"))
    index: dict[str, dict[str, Any]] = {}
    for category in data.get("categories", []):
        for item in category.get("items", []):
            name = item.get("name")
            if name:
                index[name] = item
    return index


def _download_photo_as_png(url: str, out_path: Path) -> bool:
    request = urllib.request.Request(
        _url_to_request_url(url),
        headers={
            "User-Agent": "Mozilla/5.0",
            "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
        },
    )
    with urllib.request.urlopen(request, timeout=25) as response:
        data = response.read(10_000_000)

    with Image.open(BytesIO(data)) as img:
        img = ImageOps.exif_transpose(img)
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGB")
        elif img.mode == "RGBA":
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.getchannel("A"))
            img = bg
        out_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(out_path, "PNG", optimize=True)
    return True


def _cached_original_photo(project_dir: Path, menu_index: dict[str, dict[str, Any]], dish_name: str) -> Path | None:
    if not dish_name or dish_name not in menu_index:
        return None

    cache_dir = project_dir / "original_dish_images"
    slug = _safe_filename(dish_name)
    existing = cache_dir / f"{slug}.png"
    if existing.exists():
        return existing

    photos = menu_index[dish_name].get("photos") or []
    for idx, photo in enumerate(photos[:4], start=1):
        url = photo.get("url")
        if not url:
            continue
        fallback = cache_dir / f"{slug}_{hashlib.md5(url.encode('utf-8')).hexdigest()[:8]}.png"
        try:
            if _download_photo_as_png(url, fallback):
                if not existing.exists():
                    fallback.replace(existing)
                return existing
        except Exception:
            continue
    return None


def _is_compatible_reference(new_name: str, base_name: str) -> bool:
    if not new_name or not base_name:
        return False

    if "汤" in new_name:
        return "汤" in base_name
    if "蒸" in new_name:
        return "蒸" in base_name or ("剁椒" in new_name and "剁椒" in base_name)
    if "焖" in new_name:
        return "焖" in base_name
    if "砂锅" in new_name:
        return "砂锅" in base_name
    if "干锅" in new_name:
        return "干锅" in base_name
    if "小炒" in new_name:
        return "炒" in base_name or "辣椒炒肉" in base_name
    if "擂椒" in new_name:
        return "擂" in base_name or "青椒" in base_name
    if "荷塘" in new_name:
        return "荷塘" in base_name
    if "口味" in new_name:
        return "口味" in base_name
    if "剁椒" in new_name:
        return "剁椒" in base_name
    return False


def _original_reference_for(
    project_dir: Path,
    menu_index: dict[str, dict[str, Any]],
    item: dict[str, Any],
) -> tuple[str, Path | None, bool]:
    base = str(item.get("base_menu_item") or "").strip()
    new_name = str(item.get("name") or "").strip()
    if not base or base not in menu_index or not _is_compatible_reference(new_name, base):
        return "", None, False

    image_path = _cached_original_photo(project_dir, menu_index, base)
    if not image_path:
        return "", None, False
    return base, image_path, False


def _fit_reason_for_slide(item: dict[str, Any], has_reference: bool, restaurant_name: str) -> str:
    name = str(item.get("name") or "")
    base = str(item.get("base_menu_item") or "").strip()
    if has_reference:
        return item.get("fit_reason", "")
    if "汤" in name:
        return f"{restaurant_name}现有菜单未找到直接对应的汤类原菜，本页按新增汤品机会呈现：用莲藕清甜与骨汤温润补足解辣、暖胃、家庭共享场景。"
    if base:
        category = _category_for(item)
        return f"现有菜单未找到与“{name}”直接对应的同类原菜，因此不按“{base}”做转化呈现；本页作为{category}新增合作菜品验证。"
    return item.get("fit_reason", "")


def _header_title_size(title: str) -> int:
    if len(title) > 25:
        return 20
    if len(title) > 18:
        return 22
    return 24


def _add_header(slide, prs: Presentation, title: str, section: str = "") -> int:
    width = prs.slide_width
    _add_text(
        slide,
        title,
        HEADER_LEFT,
        HEADER_TOP,
        HEADER_TITLE_WIDTH,
        HEADER_TITLE_HEIGHT,
        size=_header_title_size(title),
        color=DARK,
        bold=True,
    )
    if section:
        _add_text(
            slide,
            section,
            HEADER_LEFT,
            HEADER_SECTION_TOP,
            HEADER_SECTION_WIDTH,
            HEADER_SECTION_HEIGHT,
            size=11,
            color=RED,
            bold=True,
        )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        HEADER_LEFT,
        HEADER_RULE_TOP,
        width - HEADER_LEFT * 2,
        HEADER_RULE_HEIGHT,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    return HEADER_CONTENT_TOP


def _add_image_caption(slide, label: str, left: int, top: int, width: int, *, fill=INK, size: int = 13) -> None:
    _add_text(slide, label, left, top, width, 330_000, size=size, color=WHITE, bold=True, fill=fill, line=fill)


def _build_cover(slide, prs: Presentation, project_dir: Path, hero_item: dict[str, Any], restaurant_name: str) -> None:
    img = _image_path(project_dir, hero_item)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and ("幻师" in shape.text or "菜品推荐" in shape.text):
            _set_shape_text(shape, f"{restaurant_name} × 荷仙莲藕\n合作菜品提案", size=32, color=RED, bold=True)

    hero = _find_cover_hero_picture(slide, prs)
    if hero is None:
        raise RuntimeError("Could not find the existing right-side cover image shape to replace.")
    if not _replace_picture_content(slide, hero, img, crop_to_fill=True):
        raise RuntimeError(f"Cover hero image not found: {img}")


def _build_overview(slide, prs: Presentation, project_dir: Path, items: list[dict[str, Any]]) -> None:
    _clear_slide_for_rebuild(slide)
    content_top = _add_header(slide, prs, "合作菜品推荐总览", "第6页起：湘菜风格合作菜品方案")
    top_names = ["擂椒藕丁", "藕泥蒸排骨", "干锅藕片", "凉拌藕丁", "筒骨莲藕汤"]
    top_items = [next((x for x in items if x.get("name") == n), None) for n in top_names]
    top_items = [x for x in top_items if x]
    x0, y0 = 350_000, content_top
    gap = 55_000
    card_w = (prs.slide_width - x0 * 2 - gap * 4) // 5
    card_h = 2_680_000
    for idx, item in enumerate(top_items[:5]):
        x = x0 + idx * (card_w + gap)
        _add_picture_fill(slide, _image_path(project_dir, item), x, y0, card_w, card_h, label=item["name"])
        _add_image_caption(slide, item["name"], x + 75_000, y0 + card_h - 720_000, card_w - 150_000, fill=RED, size=14)
        _add_text(
            slide,
            _flavor_for(item["name"], item.get("method_brief", "")),
            x + 75_000,
            y0 + card_h - 360_000,
            card_w - 150_000,
            250_000,
            size=9,
            color=WHITE,
            bold=True,
            fill=INK,
            line=INK,
        )
    categories = [
        ("湘味热炒", "锅气、香辣、下饭", "擂椒藕丁 / 干锅藕片 / 莲藕小炒肉"),
        ("汤品蒸菜", "温润、平衡辣味", "藕泥蒸排骨 / 莲藕排骨汤 / 筒骨莲藕汤"),
        ("开胃凉菜", "酸辣脆爽、提升翻台前菜", "凉拌藕丁 / 酸辣藕尖"),
        ("甜品收尾", "桂花香、软糯记忆点", "桂花糯米藕"),
    ]
    for idx, (name, desc, examples) in enumerate(categories):
        x = 500_000 + (idx % 2) * 5_720_000
        y = 4_520_000 + (idx // 2) * 760_000
        _add_box(slide, x, y, 5_300_000, 560_000, fill=INK, line=INK, radius=False)
        _add_text(slide, name, x + 170_000, y + 60_000, 1_420_000, 230_000, size=14, color=GOLD, bold=True)
        _add_text(slide, desc, x + 1_650_000, y + 65_000, 3_220_000, 210_000, size=11, color=WHITE, bold=True)
        _add_text(slide, examples, x + 1_650_000, y + 300_000, 3_350_000, 190_000, size=8, color=RGBColor(235, 224, 205))


def _build_category_map(slide, prs: Presentation, project_dir: Path, items: list[dict[str, Any]], restaurant_name: str) -> None:
    _clear_slide_for_rebuild(slide)
    content_top = _add_header(slide, prs, "从“莲藕产品”到“湘菜菜单”的转译", "分类介绍合作菜品")
    groups = ["湘味热炒", "汤品蒸菜", "开胃凉菜", "香酥小吃", "甜品收尾"]
    grouped = {group: [x for x in items if _category_for(x) == group] for group in groups}
    for idx, group in enumerate(groups):
        top_row = idx < 3
        x = 430_000 + (idx % 3) * 3_820_000 if top_row else 2_360_000 + (idx - 3) * 3_820_000
        y = content_top + 20_000 if top_row else content_top + 2_360_000
        tile_w, tile_h = 3_520_000, 2_160_000
        names = " / ".join(item["name"] for item in grouped[group][:5])
        first = grouped[group][0] if grouped[group] else None
        _add_picture_fill(slide, _image_path(project_dir, first) if first else None, x, y, tile_w, tile_h, label=group)
        _add_box(slide, x, y + tile_h - 730_000, tile_w, 730_000, fill=INK, line=INK, radius=False)
        _add_text(slide, group, x + 180_000, y + tile_h - 650_000, 1_500_000, 260_000, size=17, color=GOLD, bold=True)
        _add_text(slide, names, x + 180_000, y + tile_h - 350_000, tile_w - 360_000, 260_000, size=9, color=WHITE, bold=True)
    _add_text(slide, f"保持{restaurant_name}“现炒、下饭、湘味”的菜单心智", 1_160_000, 6_180_000, 9_850_000, 330_000, size=16, color=GREEN, bold=True)


def _build_dish_slide(
    slide,
    prs: Presentation,
    project_dir: Path,
    item: dict[str, Any],
    index: int,
    menu_index: dict[str, dict[str, Any]],
    restaurant_name: str,
    support_index: int = 1,
) -> None:
    _clear_slide_for_rebuild(slide)
    name = item.get("name", f"合作菜品{index}")
    category = _category_for(item)
    method = item.get("method_brief", "")
    ref_name, ref_image, _ = _original_reference_for(project_dir, menu_index, item)
    has_reference = bool(ref_name and ref_image)
    section = f"{category} · 原菜品：{ref_name}" if has_reference else f"{category} · 新增合作菜品"
    content_top = _add_header(slide, prs, name, section)

    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = 300_000
    gap = 120_000
    pane_top = content_top
    pane_h = slide_h - pane_top - 250_000
    pane_w = (slide_w - margin * 2 - gap) // 2
    left_x = margin
    right_x = margin + pane_w + gap

    if has_reference:
        _add_picture_fill(slide, ref_image, left_x, pane_top, pane_w, pane_h, label="原菜品图片待补充")
        _add_picture_fill(slide, _image_path(project_dir, item), right_x, pane_top, pane_w, pane_h, label=name)
        _add_image_caption(slide, f"原菜品｜{ref_name}", left_x + 120_000, pane_top + 120_000, pane_w - 240_000, fill=INK, size=13)
        _add_image_caption(slide, f"合作新菜｜{name}", right_x + 120_000, pane_top + 120_000, pane_w - 240_000, fill=RED, size=13)
        band_h = 1_620_000
        band_y = slide_h - band_h - 250_000
        _add_box(slide, margin, band_y, slide_w - margin * 2, band_h, fill=INK, line=INK, radius=False)
        col_w = (slide_w - margin * 2 - 360_000) // 2
        left_text = margin + 190_000
        right_text = margin + col_w + 330_000
    else:
        hero_w = 7_040_000
        side_x = margin + hero_w + gap
        side_w = slide_w - side_x - margin
        support_h = 2_120_000
        info_y = pane_top + support_h + 120_000
        info_h = pane_h - support_h - 120_000

        _add_picture_fill(slide, _image_path(project_dir, item), margin, pane_top, hero_w, pane_h, label=name)
        _add_image_caption(slide, f"合作新菜｜{name}", margin + 120_000, pane_top + 120_000, hero_w - 240_000, fill=RED, size=13)
        support_path = _support_image_path(project_dir, support_index)
        _add_picture_fill(slide, support_path, side_x, pane_top, side_w, support_h, label="餐厅场景")
        _add_box(slide, side_x, info_y, side_w, info_h, fill=LIGHT, line=GREEN, radius=False)

        col_w = side_w - 360_000
        left_text = side_x + 180_000
        right_text = left_text
        band_y = info_y

        _add_text(slide, "新增逻辑", left_text, band_y + 130_000, 1_350_000, 220_000, size=12, color=GREEN, bold=True)
        _add_text(slide, _short(_fit_reason_for_slide(item, has_reference, restaurant_name), 76), left_text, band_y + 400_000, col_w, 390_000, size=8, color=DARK)
        _add_text(slide, "口味卖点", left_text, band_y + 850_000, 1_550_000, 220_000, size=12, color=GREEN, bold=True)
        _add_text(slide, _flavor_for(name, method), left_text, band_y + 1_110_000, col_w, 210_000, size=9, color=DARK, bold=True)
        _add_text(slide, "制作方法", left_text, band_y + 1_390_000, 1_550_000, 220_000, size=12, color=GREEN, bold=True)
        _add_method_steps(slide, method, left_text, band_y + 1_650_000, col_w, 520_000, color=DARK, size=8)
        _add_text(slide, "莲藕角色", left_text, band_y + 2_230_000, 1_350_000, 220_000, size=12, color=GREEN, bold=True)
        _add_text(slide, _short(item.get("lotus_role", ""), 48), left_text, band_y + 2_490_000, col_w, 260_000, size=8, color=MUTED)
        _add_text(slide, f"合作菜品 {index:02d}", 10_000_000, 620_000, 1_300_000, 260_000, size=10, color=GOLD, bold=True)
        return

    _add_text(slide, "转译逻辑", left_text, band_y + 120_000, 1_350_000, 230_000, size=12, color=GOLD, bold=True)
    _add_text(slide, _short(_fit_reason_for_slide(item, has_reference, restaurant_name), 84), left_text, band_y + 370_000, col_w - 140_000, 340_000, size=8, color=WHITE)
    _add_text(slide, "莲藕角色", left_text, band_y + 760_000, 1_350_000, 230_000, size=12, color=GOLD, bold=True)
    _add_text(slide, _short(item.get("lotus_role", ""), 66), left_text, band_y + 1_020_000, col_w - 140_000, 350_000, size=8, color=RGBColor(235, 224, 205))

    _add_text(slide, "口味卖点", right_text, band_y + 120_000, 1_350_000, 230_000, size=12, color=GOLD, bold=True)
    _add_text(slide, _flavor_for(name, method), right_text, band_y + 395_000, col_w - 140_000, 260_000, size=12, color=WHITE, bold=True)
    _add_text(slide, "制作方法", right_text, band_y + 680_000, 1_350_000, 230_000, size=12, color=GOLD, bold=True)
    _add_method_steps(slide, method, right_text, band_y + 940_000, col_w - 140_000, 560_000, color=RGBColor(235, 224, 205), size=8)
    _add_text(slide, f"合作菜品 {index:02d}", 10_000_000, 620_000, 1_300_000, 260_000, size=10, color=GOLD, bold=True)


def _build_launch_plan(slide, prs: Presentation, project_dir: Path, items: list[dict[str, Any]]) -> None:
    _clear_slide_for_rebuild(slide)
    content_top = _add_header(slide, prs, "建议试点组合：先用高确定性菜品打样", "上新节奏建议")
    steps = [
        ("第一轮试菜", "擂椒藕丁 / 藕泥蒸排骨 / 干锅藕片", "验证口味接受度与出餐效率"),
        ("第二轮菜单化", "凉拌藕丁 / 筒骨莲藕汤 / 桂花糯米藕", "补齐凉菜、汤品、甜品场景"),
        ("第三轮区域试点", "选择 3-5 家高客流门店", "观察点单率、复购反馈、桌均带动"),
    ]
    for idx, (title, dishes, note) in enumerate(steps):
        y = content_top + idx * 1_350_000
        _add_box(slide, 760_000, y, 5_000_000, 920_000, fill=LIGHT)
        _add_text(slide, title, 960_000, y + 130_000, 1_400_000, 280_000, size=17, color=RED, bold=True)
        _add_text(slide, dishes, 2_420_000, y + 130_000, 2_900_000, 250_000, size=13, color=DARK, bold=True)
        _add_text(slide, note, 2_420_000, y + 460_000, 2_900_000, 250_000, size=10, color=MUTED)
    for idx, item in enumerate(items[:4]):
        _add_picture_fit(slide, _image_path(project_dir, item), 6_270_000 + (idx % 2) * 2_550_000, content_top + (idx // 2) * 2_220_000, 2_160_000, 2_160_000)


def _build_supply_slide(slide, prs: Presentation, restaurant_name: str) -> None:
    _clear_slide_for_rebuild(slide)
    content_top = _add_header(slide, prs, "供应与出餐适配：把莲藕变成门店可执行菜品", "落地重点")
    rows = [
        ("原料形态", "鲜切藕丁、藕片、藕带、藕泥、藕夹半成品", "按菜型定制规格，减少门店前处理"),
        ("出餐逻辑", "热炒快出、凉菜预拌、汤品预炖、甜品预制", f"匹配{restaurant_name}高峰翻台节奏"),
        ("口味方向", "香辣下饭 + 清爽解腻 + 温润汤品", "补足现有菜单的桌搭层次"),
        ("价格状态", "当前建议价格均为待核实", "需结合门店成本、规格、毛利目标测算"),
    ]
    y = content_top
    for title, content, note in rows:
        _add_box(slide, 860_000, y, 10_200_000, 760_000, fill=LIGHT)
        _add_text(slide, title, 1_050_000, y + 150_000, 1_500_000, 280_000, size=15, color=RED, bold=True)
        _add_text(slide, content, 2_630_000, y + 120_000, 3_850_000, 300_000, size=13, color=DARK, bold=True)
        _add_text(slide, note, 6_700_000, y + 130_000, 3_900_000, 300_000, size=11, color=MUTED)
        y += 980_000


def _build_closing(slide, prs: Presentation, project_dir: Path, items: list[dict[str, Any]], restaurant_name: str) -> None:
    _clear_slide_for_rebuild(slide)
    width, height = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()
    _add_text(slide, f"期待与{restaurant_name}共创莲藕湘菜新爆品", 900_000, 1_000_000, 6_200_000, 620_000, size=30, color=RED, bold=True)
    _add_text(slide, "从热炒、凉菜、蒸菜、汤品到甜品，用莲藕建立更完整的餐桌记忆点。", 920_000, 1_850_000, 6_000_000, 560_000, size=16, color=DARK)
    _add_text(slide, "下一步：试菜打样 · 成本核算 · 门店试点 · 菜单上新", 920_000, 2_650_000, 6_200_000, 380_000, size=14, color=GREEN, bold=True)
    for idx, item in enumerate(items[:3]):
        _add_picture_fit(slide, _image_path(project_dir, item), 7_450_000 + idx * 1_400_000, 1_280_000 + idx * 1_100_000, 1_760_000, 1_760_000)
    _add_text(slide, "谢谢观看", 930_000, 5_720_000, 2_100_000, 360_000, size=20, color=DARK, bold=True)


def create_ppt(
    template_path: Path,
    project_dir: Path,
    output_name: str = "",
    cover_dish_name: str = "",
    restaurant_name: str = "",
) -> Path:
    prs = Presentation(str(template_path))
    opportunity_data = _load_opportunity_data(project_dir)
    restaurant_name = _restaurant_name_from_project(project_dir, opportunity_data, restaurant_name)
    items = opportunity_data.get("opportunities", [])
    menu_index = _load_menu_photo_index(project_dir)
    if len(items) < 17:
        raise RuntimeError("lotus_root_opportunities.json should contain at least 17 opportunity items.")

    _build_cover(prs.slides[0], prs, project_dir, _select_cover_item(items, project_dir, cover_dish_name), restaurant_name)
    _replace_text_contains(
        prs,
        2,
        {
            "三、香脆藕条介绍": f"三、{restaurant_name}菜单机会诊断",
            "四、四类藕产品推介": "四、合作菜品推荐方案",
        },
    )

    _build_overview(prs.slides[5], prs, project_dir, items)
    _build_category_map(prs.slides[6], prs, project_dir, items, restaurant_name)
    support_index = 0
    for offset, item in enumerate(items, start=8):
        ref_name, ref_image, _ = _original_reference_for(project_dir, menu_index, item)
        if not (ref_name and ref_image):
            support_index += 1
        _build_dish_slide(
            prs.slides[offset - 1],
            prs,
            project_dir,
            item,
            offset - 7,
            menu_index,
            restaurant_name,
            support_index=max(support_index, 1),
        )

    # Keep the source deck's original closing page. The two template pages before it
    # are not needed for this proposal, so remove them after custom dish pages are built.
    if len(prs.slides) >= 27:
        _delete_slide(prs, 25)
        _delete_slide(prs, 24)

    if not output_name:
        output_name = f"{_safe_filename(restaurant_name)}_lotus_coop_dish_proposal.pptx"
    out = project_dir / output_name
    try:
        prs.save(str(out))
    except PermissionError:
        fallback = out.with_name(f"{out.stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{out.suffix}")
        prs.save(str(fallback))
        return fallback
    return out


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    output = create_ppt(Path(args.template), Path(args.project_dir), args.output, args.cover_dish_name, args.restaurant_name)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
